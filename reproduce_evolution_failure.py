import sys
import os
import io
from pptx import Presentation

# Setup imports
sys.path.append(os.getcwd())
sys.path.append(os.path.join(os.getcwd(), 'api'))

try:
    from api.ppt_shared import add_chart_slide, create_logo_slide, create_intro_slide
except ImportError:
    # Handle if running from different dir
    sys.path.append(os.path.join(os.getcwd(), 'api'))
    from ppt_shared import add_chart_slide, create_logo_slide, create_intro_slide

def create_test_ppt():
    prs = Presentation()
    create_logo_slide(prs)
    create_intro_slide(prs, "Test Evolución", "01/02/2026")
    
    # Simulate Evolution Data (Date strings as keys?)
    # Usually the dashboard sends: [{'Date': '2025-01', 'Toyota': 100}, {'Date': '2025-02', 'Toyota': 102}]
    evolution_data = [
        {'Fecha': '2025-01-01', 'Toyota RAV4': 25000000, 'Mazda CX-5': 24500000},
        {'Fecha': '2025-02-01', 'Toyota RAV4': 25200000, 'Mazda CX-5': 24600000},
        {'Fecha': '2025-03-01', 'Toyota RAV4': 25500000, 'Mazda CX-5': 24800000},
        {'Fecha': '2025-04-01', 'Toyota RAV4': 26000000, 'Mazda CX-5': 25000000},
    ]
    
    chart_info = {
        'chart_title': 'Evolución de Precios',
        'chart_type': 'line',
        'data': evolution_data
    }
    
    print("Adding Evolution Chart...")
    add_chart_slide(prs, chart_info, currency_symbol='$')
    
    # Simulate Compare Data (Categories)
    compare_data = [
        {'Versión': 'RAV4 LE 2.0', 'Precio Lista': 25000000, 'Bono': 1000000},
        {'Versión': 'RAV4 XLE 2.5', 'Precio Lista': 29000000, 'Bono': 1500000},
    ]
    
    chart_info_compare = {
        'chart_title': 'Comparación de Versiones',
        'chart_type': 'bar', # Front might send 'bar' for column chart look
        'data': compare_data
    }
    
    print("Adding Compare Chart...")
    add_chart_slide(prs, chart_info_compare, currency_symbol='$')
    
    filename = 'test_evolution_debug.pptx'
    prs.save(filename)
    print(f"Saved to {filename}")

if __name__ == '__main__':
    try:
        create_test_ppt()
    except Exception as e:
        print(f"CRITICAL ERROR: {e}")
        import traceback
        traceback.print_exc()
