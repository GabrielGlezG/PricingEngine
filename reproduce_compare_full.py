import sys
import os
import json
from datetime import datetime
from pptx import Presentation

# Setup imports
sys.path.append(os.getcwd())
sys.path.append(os.path.join(os.getcwd(), 'api'))

try:
    from api.ppt_shared import add_chart_slide, add_table_slide, create_logo_slide, create_intro_slide
except ImportError:
    sys.path.append(os.path.join(os.getcwd(), 'api'))
    from ppt_shared import add_chart_slide, add_table_slide, create_logo_slide, create_intro_slide

def create_compare_ppt():
    prs = Presentation()
    create_logo_slide(prs)
    create_intro_slide(prs, "Comparación de Vehículos", "01/02/2026")
    
    # Simulate Summary Table Data
    summary_data = [
        {'Marca': 'Toyota', 'Modelo': 'RAV4', 'Versión': 'LE 2.0', 'Segmento': 'SUV', 'Precio Actual': 25000000},
        {'Marca': 'Toyota', 'Modelo': 'RAV4', 'Versión': 'XLE 2.5', 'Segmento': 'SUV', 'Precio Actual': 29000000},
        {'Marca': 'Mazda', 'Modelo': 'CX-5', 'Versión': 'Sport', 'Segmento': 'SUV', 'Precio Actual': 26000000},
    ]
    
    print("Adding Summary Table...")
    add_table_slide(prs, 'Resumen Comparación', summary_data, '$')
    
    # Simulate Evolution Data with Long Names (Version Specific)
    # Gaps included to test smooth lines
    evolution_data = [
        {'Fecha': '2025-01-01', 'Toyota RAV4 LE 2.0': 24000000, 'Toyota RAV4 XLE 2.5': 28000000, 'Mazda CX-5 Sport': 25500000},
        {'Fecha': '2025-02-01', 'Toyota RAV4 LE 2.0': 0,        'Toyota RAV4 XLE 2.5': 28500000, 'Mazda CX-5 Sport': 25800000}, # Gap for LE
        {'Fecha': '2025-03-01', 'Toyota RAV4 LE 2.0': 25000000, 'Toyota RAV4 XLE 2.5': 0,        'Mazda CX-5 Sport': 26000000}, # Gap for XLE
        {'Fecha': '2025-04-01', 'Toyota RAV4 LE 2.0': 25000000, 'Toyota RAV4 XLE 2.5': 29000000, 'Mazda CX-5 Sport': 0},        # Gap for Mazda
    ]
    
    chart_info = {
        'chart_title': 'Evolución de Precios',
        'chart_type': 'line',
        'data': evolution_data
    }
    
    print("Adding Evolution Chart...")
    add_chart_slide(prs, chart_info, currency_symbol='$')
    
    filename = 'test_compare_full.pptx'
    prs.save(filename)
    print(f"Saved to {filename}")

if __name__ == '__main__':
    try:
        create_compare_ppt()
    except Exception as e:
        print(f"CRITICAL ERROR: {e}")
        import traceback
        traceback.print_exc()
