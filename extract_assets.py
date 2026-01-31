from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

def extract_images(path):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return

    prs = Presentation(path)
    output_dir = "public/extracted"
    os.makedirs(output_dir, exist_ok=True)
    
    print(f"Extracting images from: {path}")
    
    image_count = 0
    for i, slide in enumerate(prs.slides):
        print(f"Scanning Slide {i+1}...")
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                filename = f"slide_{i+1}_image_{image_count}.{ext}"
                full_path = os.path.join(output_dir, filename)
                with open(full_path, 'wb') as f:
                    f.write(image.blob)
                print(f"  Saved: {filename}")
                image_count += 1
            
            # Check for Background Images (if accessible via shapes, often simple BGs are not shapes)
            # Complex BGs are often locked pictures.

extract_images("public/FORMATO MAURICIO.pptx")
