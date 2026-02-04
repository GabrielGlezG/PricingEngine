import os
import sys
import base64
import io
import importlib.util
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.chart import (
    XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION, 
    XL_TICK_MARK, XL_MARKER_STYLE, XL_TICK_LABEL_POSITION
)
from pptx.oxml.ns import qn

# --- BRAND COLORS (Institutional) ---
DARK_BLUE = RGBColor(30, 41, 59)  # Slate 900 #1E293B
DEEP_NAVY = RGBColor(13, 40, 65)  # #0D2841 (Cover BG)
LIGHT_BLUE = RGBColor(71, 85, 105) # Slate 600
WHITE = RGBColor(255, 255, 255)

# --- ASSET LOADING ---
# Robust asset loading bypassing path issues
try:
    # Try to find ppt_assets.py in the same directory as this file
    current_dir = os.path.dirname(os.path.abspath(__file__))
    assets_path = os.path.join(current_dir, 'ppt_assets.py')
    
    if os.path.exists(assets_path):
        spec = importlib.util.spec_from_file_location("ppt_assets", assets_path)
        ppt_assets = importlib.util.module_from_spec(spec)
        spec.loader.exec_module(ppt_assets)
        LOGO_B64 = getattr(ppt_assets, 'LOGO_B64', None)
        BG_B64 = getattr(ppt_assets, 'BG_B64', None)
    else:
        # Fallback if running from root
        try:
            from api.ppt_assets import LOGO_B64, BG_B64
        except:
            LOGO_B64 = None
            BG_B64 = None
except Exception as e:
    print(f"Warning: Failed to load assets: {e}")
    LOGO_B64 = None
    BG_B64 = None

# --- HELPER FUNCTIONS ---

def format_value(val, fmt=None, currency='$'):
    if val is None: return "-"
    try:
        val_float = float(val)
        if fmt == 'currency':
             if currency == 'UF':
                 return f"UF {val_float:,.2f}".replace(",", ".")
             else:
                 return f"{currency} {val_float:,.0f}".replace(",", ".")
        elif fmt == 'percent':
            return f"{val_float * 100:.1f}%"
        elif fmt == 'integer':
             return f"{val_float:,.0f}".replace(",", ".")
        elif isinstance(val, (int, float)):
            return f"{val_float:,.0f}" if val_float.is_integer() else f"{val_float:.2f}"
    except:
        pass
    return str(val)

def set_font(shape, font_name="Avenir Medium", font_size=None, bold=False, color=None):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = font_name
            if font_size: run.font.size = font_size
            if bold is not None: run.font.bold = bold
            if color: run.font.color.rgb = color

def get_image_stream(b64_string):
    if not b64_string:
        return None
    return io.BytesIO(base64.b64decode(b64_string))

# --- COMMON SLIDE CREATORS ---

def create_logo_slide(prs):
    """Creates a slide with a large centered logo on Deep Navy background using embedded asset."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = DEEP_NAVY # #0D2841
    
    try:
        img_stream = get_image_stream(LOGO_B64)
        if img_stream:
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            logo_width = Inches(5.5) 
            left = (slide_width - logo_width) / 2
            
            pic = slide.shapes.add_picture(img_stream, left, 0, width=logo_width)
            
            top = (slide_height - pic.height) / 2
            pic.top = int(top)
        else:
            print("Warning: LOGO_B64 not available")
    except Exception as e:
        print(f"Error adding logo to slide: {e}")

def create_intro_slide(prs, title, date_str):
    """Creates intro slide with split background using embedded asset."""
    slide = prs.slides.add_slide(prs.slide_layouts[6]) 
    
    try:
        img_stream = get_image_stream(BG_B64)
        if img_stream:
            slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
    except Exception as e:
        print(f"Error adding background to intro: {e}")
        
    # Text Layout for Split Background (Left side white space)
    left = Inches(0.8) # Left padding
    top = Inches(2.8)  # Vertical center alignment roughly
    width = Inches(4.5) 
    height = Inches(2.5)
    
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = title # Keep original casing or .title()
    p_title.font.name = "Avenir Black"
    p_title.font.size = Pt(40) # Large Title
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(0, 0, 0) # Black text on White part
    p_title.alignment = PP_ALIGN.LEFT
    
    p_date = tf.add_paragraph()
    p_date.text = f"\nGenerado: {date_str}"
    # Explicitly style the run to avoid inheritance issues
    if p_date.runs:
        run = p_date.runs[0]
        run.font.name = "Avenir Medium"
        run.font.size = Pt(16)
        run.font.color.rgb = LIGHT_BLUE
    p_date.alignment = PP_ALIGN.LEFT

# --- SHARED CHART & TABLE LOGIC ---

def add_table_slide(prs, title, rows, currency_symbol='$'):
    if not rows: return
    
    # Configuration
    MAX_ROWS = 12
    MAX_DATA_COLS = 7 # Reduced from 10 to ensure wide columns for currency (No wrapping)
    
    # 1. Prepare Columns (Headers)
    all_headers = list(rows[0].keys())
    fixed_header = all_headers[0] # Assume 'Fecha' or 'Marca' is first
    data_headers = all_headers[1:]
    
    # Chunk Columns (Horizontal Split) - Outer Loop (Group by Series/Topic)
    col_chunks = [data_headers[i:i + MAX_DATA_COLS] for i in range(0, len(data_headers), MAX_DATA_COLS)]
    
    slide_count = 0
    
    for c_idx, col_chunk in enumerate(col_chunks):
        # Current subset of headers for this "Sheet"
        current_headers = [fixed_header] + col_chunk
        
        # 2. Chunk Rows (Vertical Split) - Inner Loop (Pagination over time/items)
        row_chunks = [rows[i:i + MAX_ROWS] for i in range(0, len(rows), MAX_ROWS)]
        
        for r_idx, row_chunk in enumerate(row_chunks):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            
            # Smart Title: "Title (Part X)" or simple Continuation
            # "Evolución de Precios" -> "Evolución de Precios (Cont.)"
            suffix = ""
            if slide_count > 0: suffix = " (Cont.)"
            slide.shapes.title.text = f"{title}{suffix}"
            
            set_font(slide.shapes.title, font_name="Avenir Black", font_size=Pt(28), bold=True, color=DARK_BLUE)
            
            # Create Table
            rows_num = len(row_chunk) + 1
            cols_num = len(current_headers)
            
            # Dynamic Width Calculation? Standard is 9 inches total.
            # If fewer columns, use full width? Yes.
            
            # 16:9 Layout Adjustments (Width 13.33")
            # Centered Wider Table: Width 12", Margin (13.33 - 12)/2 = 0.665"
            shape = slide.shapes.add_table(rows_num, cols_num, Inches(0.665), Inches(1.5), Inches(12), Inches(0.4 * rows_num))
            table = shape.table
            
            # Render Header Row
            for c, header in enumerate(current_headers):
                cell = table.cell(0, c)
                cell.text = str(header)
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE 
                tf = cell.text_frame.paragraphs[0]
                tf.font.color.rgb = WHITE
                tf.font.bold = True
                tf.font.name = "Avenir Medium"
                tf.font.size = Pt(10)
                tf.alignment = PP_ALIGN.CENTER

            # Render Data Rows
            for r, row_data in enumerate(row_chunk):
                for c, header in enumerate(current_headers):
                    val = row_data.get(header)
                    cell = table.cell(r + 1, c)
                    
                    h_str = str(header).lower()
                    t_str = str(title).lower()
                    fmt = None
                    
                    if any(x in h_str for x in ['%', 'percent', 'variación', 'variation', 'coef', 'descuento', 'volatilidad']):
                        fmt = 'percent'
                    elif any(x in h_str for x in ['cantidad', 'cant.', 'volumen', 'versiones', 'total', 'numero', 'count']):
                        fmt = 'integer'
                    elif any(x in h_str for x in ['precio', 'price', 'monto', 'valor', 'bono', 'lista', 'costo', 'avg', 'min', 'max', 'promedio']):
                        fmt = 'currency'
                    elif any(x in t_str for x in ['volatilidad', 'volatility', 'variación', 'variation', 'share', 'participación', 'discount', 'descuento']):
                         if isinstance(val, (int, float)) and not any(x in h_str for x in ['fecha', 'date', 'year', 'año', 'mes']): 
                            fmt = 'percent'
                    elif "precio" in t_str or "price" in t_str:
                         if isinstance(val, (int, float)): fmt = 'currency'
                    
                    if 'año' in h_str or 'year' in h_str: fmt = None

                    cell.text = format_value(val, fmt, currency_symbol)
                    tf = cell.text_frame.paragraphs[0]
                    tf.font.name = "Avenir Medium"
                    tf.font.size = Pt(9)
                    
                    if fmt in ['currency', 'percent', 'integer']:
                         tf.alignment = PP_ALIGN.RIGHT
                    else:
                         tf.alignment = PP_ALIGN.LEFT
            
            slide_count += 1

def add_chart_slide(prs, chart_info, currency_symbol='$'):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = chart_info.get('chart_title', 'Gráfico')
    set_font(slide.shapes.title, font_name="Avenir Black", font_size=Pt(28), bold=True, color=DARK_BLUE)
    
    chart_type = chart_info.get('chart_type', 'bar')
    rows = chart_info.get('data', [])
    if not rows: return

    headers = list(rows[0].keys())
    categories = [str(r[headers[0]]) for r in rows]
    series_names = headers[1:] # Assuming rest are series
    
    
    # Heuristics based on chart name/title to apply specific formatting
    name_lower = str(chart_info.get('name') or chart_info.get('chart_title') or '').lower()
    
    # 1. Chart Types & Data Preparation
    if chart_type == 'scatter':
         ppt_chart_type = XL_CHART_TYPE.BUBBLE
         chart_data = BubbleChartData() 
         
         for r in rows:
            label = str(r[headers[0]])
            try:
                x_key = next((k for k in headers if 'volumen' in k.lower() or 'cant' in k.lower()), headers[1])
                y_key = next((k for k in headers if 'precio' in k.lower() or 'price' in k.lower()), headers[2])
                x_val = float(r.get(x_key, 0))
                y_val = float(r.get(y_key, 0))
                size_raw = float(r.get(x_key, 0))
                size_val = abs(size_raw) if abs(size_raw) > 0.0 else 1.0
                series = chart_data.add_series(label)
                series.add_data_point(x_val, y_val, size_val)
            except:
                continue
    else:
        ppt_chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
        if chart_type == 'line': ppt_chart_type = XL_CHART_TYPE.LINE
        elif chart_type == 'stacked': ppt_chart_type = XL_CHART_TYPE.COLUMN_STACKED
        
        is_evolution = 'evolución' in name_lower or 'evolution' in name_lower
        
        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s_name in series_names:
            values = []
            for r in rows:
                val = r.get(s_name, 0)
                try:
                    # Logic: If Evolution, treat 0 as None (Gap)
                    # Otherwise default to 0.0
                    fval = float(val) if val is not None else 0.0
                    
                    if is_evolution and fval == 0.0:
                        values.append(None)
                    else:
                        values.append(fval)
                except: 
                    values.append(0.0)
            chart_data.add_series(s_name, values)

            chart_data.add_series(s_name, values)

    # 16:9 Layout Adjustments (Width 13.33")
    # Centered Wider Chart: Width 12", Margin (13.33 - 12)/2 = 0.665"
    x, y, cx, cy = Inches(0.665), Inches(1.3), Inches(12), Inches(6.0)
    graphic_frame = slide.shapes.add_chart(ppt_chart_type, x, y, cx, cy, chart_data)
    chart = graphic_frame.chart
    
    # --- CHART FORMATTING APPLIED ---
    try:
        chart.has_legend = True
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.name = "Avenir Medium"
        chart.legend.font.size = Pt(9)
        
        # --- GLOBAL AXIS FONT STYLING ---
        # Apply standard branding to axes if they validly exist
        try:
            if chart.value_axis:
                # Tick Labels
                chart.value_axis.tick_labels.font.name = "Avenir Medium"
                chart.value_axis.tick_labels.font.size = Pt(9)
                # Axis Title (e.g. "Valor")
                if chart.value_axis.has_title:
                     chart.value_axis.axis_title.text_frame.paragraphs[0].font.name = "Avenir Medium"
                     chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except: pass 

        try:
            if chart.category_axis:
                # Tick Labels
                chart.category_axis.tick_labels.font.name = "Avenir Medium"
                chart.category_axis.tick_labels.font.size = Pt(9)
                # Axis Title
                if chart.category_axis.has_title:
                     chart.category_axis.axis_title.text_frame.paragraphs[0].font.name = "Avenir Medium"
                     chart.category_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
        except: pass
        
        plot = chart.plots[0]

        # --- GLOBAL DATA LABEL FONT STYLING ---
        try:
            if plot.has_data_labels:
                dl = plot.data_labels
                dl.font.name = "Avenir Medium"
                dl.font.size = Pt(9)
        except: pass
        
        # --- NEGATIVE VALUE INVERSION (RED BARS) ---
        # Explicitly coloring points Red because invert_if_negative defaults to White/Theme
        if 'tendencia' in name_lower or 'variación' in name_lower or 'variacion' in name_lower:
             try:
                 for series in chart.series:
                     # FIRST: Disable automatic inversion
                     series.invert_if_negative = False
                     
                     points = series.points
                     # python-pptx series.values is a tuple of values
                     for i, val in enumerate(series.values):
                         if val is not None and isinstance(val, (int, float)) and val < 0:
                             try:
                                 pt = series.points[i]
                                 
                                 # Local import to ensure availability
                                 from pptx.dml.color import RGBColor
                                 
                                 # 1. Fill Red (Native API)
                                 pt.format.fill.solid()
                                 pt.format.fill.fore_color.rgb = RGBColor(255, 0, 0)
                                 
                                 # 2. Border Red
                                 pt.format.line.fill.solid() 
                                 pt.format.line.color.rgb = RGBColor(255, 0, 0)
                                 pt.format.line.width = Pt(0.75)

                             except Exception as e_pt:
                                 # Log individual point failure
                                 print(f"Failed to color point {i}: {e_pt}")

             except Exception as e:
                 print(f"Error coloring negative points: {e}")
        
        # 0. "Evolución" (Evolution) -> Line Chart, No Data Labels (Clean), Currency Axis
        if 'evolución' in name_lower or 'evolution' in name_lower:
             # Connect data points with line (don't leave gaps) using XML
             # chart.display_blanks_as = XL_DISPLAY_BLANKS_AS.SPAN (Not available in this version)
             try:
                 c_chart = chart._chartSpace.chart
                 # Ensure dispBlanksAs element exists
                 if c_chart.dispBlanksAs is None:
                    c_chart.add_dispBlanksAs()
                 
                 # Access the element directly to avoid proxy issues
                 # In some versions, it might be an oxml element directly or a proxy.
                 # Try setting the attribute on the element.
                 c_chart.dispBlanksAs.set('val', 'span')
             except Exception as e:
                 print(f"Error setting display_blanks_as: {e}")

             # Manual Layout to reserve space for Dates (Bottom) and Legend (Below Dates)
             try:
                 chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                 chart.legend.include_in_layout = False # Overlay mode (we create space manually)
                 chart.legend.font.size = Pt(9)

                 # Shrink Plot Area to 75% Height to make room at bottom
                 plot = chart.plots[0]
                 playout = plot.layout
                 playout.manual_layout.height = 0.75
                 playout.manual_layout.y = 0.0 # Top aligned
                 # playout.manual_layout.width = 1.0
                 # playout.manual_layout.x = 0.0
             except Exception as e:
                 print(f"Error setting manual layout: {e}")
             
             # Y-Axis Currency
             if chart.value_axis:
                 chart.value_axis.tick_labels.number_format = f'"{currency_symbol}" #,##0'
                 chart.value_axis.major_unit = None # Auto scale
                 chart.value_axis.tick_labels.font.name = "Avenir Medium"
                 chart.value_axis.tick_labels.font.size = Pt(9)
                 
             # Smooth Lines with Markers
             try:
                 from pptx.enum.chart import XL_MARKER_STYLE
                 for series in chart.series:
                     series.smooth = True
                     series.format.line.width = Pt(2.5)
                     # Add Markers (Dots) to match platform
                     series.marker.style = XL_MARKER_STYLE.CIRCLE
                     series.marker.size = 7
             except Exception as e:
                 print(f"Error setting markers: {e}")
             
             # X-Axis Dates (Rotate if needed)
             try:
                 category_axis = chart.category_axis
                 category_axis.tick_labels.font.name = "Avenir Medium"
                 category_axis.tick_labels.font.size = Pt(9)
                 category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
                 
                 # XML Hack for Rotation (-45 or -90 degrees)
                 txPr = category_axis._element.get_or_add_txPr()
                 # Check if bodyPr exists
                 if hasattr(txPr, 'bodyPr'):
                     bodyPr = txPr.bodyPr
                 else:
                     # Fallback for OXML element manipulation
                     from pptx.oxml.ns import qn
                     bodyPr = txPr.find(qn('a:bodyPr'))
                     if bodyPr is None:
                         bodyPr = txPr.add_bodyPr()

                 bodyPr.set('rot', '-2700000') # -45 degrees roughly
                 bodyPr.set('vert', 'horz')
             except Exception as e:
                 print(f"Error formatting evolution axis: {e}")

        # 1. "Composición" -> Integers, Vary Colors
        elif 'composición' in name_lower or 'composition' in name_lower:
            plot.vary_by_categories = True
            
            # Y-Axis Integer
            if chart.value_axis:
                chart.value_axis.tick_labels.number_format = '0'
                
                # Calculate max value to decide on major unit
                try:
                    max_val = 0
                    if chart_type == 'stacked':
                        # For stacked, max is the max sum of a row (excluding headers)
                        relevant_headers = [h for h in rows[0].keys() if h != headers[0]]
                        for r in rows:
                            row_sum = sum(float(r.get(h, 0)) for h in relevant_headers)
                            if row_sum > max_val: max_val = row_sum
                    else:
                        # For clustered, max is the single highest value
                        relevant_headers = [h for h in rows[0].keys() if h != headers[0]]
                        for r in rows:
                            for h in relevant_headers:
                                val = float(r.get(h, 0))
                                if val > max_val: max_val = val
                    
                    # Only force unit 1.0 if max value is small (e.g. <= 10) to avoid crowding
                    if max_val <= 10:
                        chart.value_axis.major_unit = 1.0
                    else:
                        chart.value_axis.major_unit = None # Auto
                except:
                    chart.value_axis.major_unit = 1.0 # Fallback
                
        # 2. "Precios" or "Estructura" or "Segmento" -> Currency, Data Labels Vertical Inside (Rotated)
        elif any(x in name_lower for x in ['precio', 'price', 'estructura', 'structure', 'segmento', 'segment']):
            # Y-Axis Currency
            if chart.value_axis:
                chart.value_axis.tick_labels.number_format = f'"{currency_symbol}" #,##0'
                # Only apply 5M steps if currency is CLP ($)
                if currency_symbol == '$':
                    chart.value_axis.major_unit = 5000000 
                else:
                     chart.value_axis.major_unit = None # Auto scale
            
            # Data Labels
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.name = "Avenir Medium"
            data_labels.font.size = Pt(8)
            data_labels.position = XL_LABEL_POSITION.INSIDE_END
            data_labels.font.color.rgb = WHITE
            data_labels.number_format = f'"{currency_symbol}" #,##0' 
            
            # --- XML HACK FOR VERTICAL TEXT (-270 deg / -90 deg visual) ---
            try:
                txPr = data_labels._element.get_or_add_txPr()
                # bodyPr usually exists on txPr creation, but let's be safe
                from pptx.oxml.ns import qn
                bodyPr = txPr.find(qn('a:bodyPr'))
                if bodyPr is None:
                     bodyPr = txPr.add_bodyPr()
                
                # Set rotation to -5400000 EMUs (approx -90 degrees)
                bodyPr.set('rot', '-5400000')
                bodyPr.set('vert', 'horz') 
            except Exception as e:
                print(f"Error rotating labels: {e}")
        
        # 2.1 "Benchmarking" -> Markers on Lines
        elif 'benchmarking' in name_lower:
             # Y-Axis Currency
             if chart.value_axis:
                 chart.value_axis.tick_labels.number_format = f'"{currency_symbol}" #,##0'
             
             # Smooth Lines with Markers
             try:
                 from pptx.enum.chart import XL_MARKER_STYLE
                 for series in chart.series:
                     series.smooth = True
                     series.format.line.width = Pt(2.5)
                     series.marker.style = XL_MARKER_STYLE.CIRCLE
                     series.marker.size = 7
             except Exception as e:
                 print(f"Error setting markers for benchmarking: {e}")
            
        # 3. "Tendencia" (Trend) -> Percent Axis, Colored Bars
        elif 'tendencia' in name_lower or 'trend' in name_lower:
            plot.vary_by_categories = True
            if chart.value_axis:
                chart.value_axis.tick_labels.number_format = '0%'
            
            # Data Labels
            plot.has_data_labels = True
            data_labels = plot.data_labels
            data_labels.font.name = "Avenir Medium"
            data_labels.font.size = Pt(8)
            data_labels.position = XL_LABEL_POSITION.INSIDE_END
            data_labels.number_format = '0%'
            data_labels.font.color.rgb = WHITE # Contrast for colored bars
            
            # Ensure negative bars are colored (not white)
            for series in chart.series:
                series.invert_if_negative = False
                
        # 4. "Volatilidad" (Volatility) -> Percent Axis, Smoothed Lines, Vertical Dates
        elif 'volatilidad' in name_lower or 'volatility' in name_lower:
            if chart.value_axis:
                chart.value_axis.tick_labels.number_format = '0.0%'
            
            # Smooth Lines
            for series in chart.series:
                series.smooth = True
                
            # Vertical X-Axis Labels (Dates)
            try:
                category_axis = chart.category_axis
                category_axis.tick_labels.font.size = Pt(9)
                category_axis.tick_label_position = XL_TICK_LABEL_POSITION.LOW
                
                # XML Hack for Rotation (-90 degrees)
                txPr = category_axis._element.get_or_add_txPr()
                bodyPr = txPr.get_or_add_bodyPr()
                bodyPr.set('rot', '-5400000')
                bodyPr.set('vert', 'horz')
            except Exception as e:
                print(f"Error formatting volatility axis: {e}")

        # 5. "Matriz" (Scatter/Bubble) -> Fix Axes
        elif 'matriz' in name_lower or chart_type == 'scatter':
             # Y-Axis (Price) -> Currency
             if chart.value_axis:
                 chart.value_axis.tick_labels.number_format = f'"{currency_symbol}" #,##0'
                 chart.value_axis.minimum_scale = 0 # Force 0 start to avoid negative values
                 if currency_symbol == '$':
                     chart.value_axis.major_unit = 5000000
             
             # X-Axis (Volume) -> Integer, Min 0
             try:
                 x_axis = chart.category_axis
                 x_axis.tick_labels.number_format = '0'
                 x_axis.minimum_scale = 0
                 x_axis.major_unit = 1
             except:
                 pass
                 
             # Reduce Bubble Size (XML Hack)
             try:
                 # chart.plots[0] is the BubblePlot
                 plot = chart.plots[0]
                 # Access the c:bubbleChart element
                 bubbleChart = plot._element
                 # Set scale to 60% (default is 100)
                 bubbleScale = bubbleChart.get_or_add_bubbleScale()
                 bubbleScale.val = 60
             except Exception as e:
                 print(f"Error scaling bubbles: {e}")

        # 6. "Benchmarking" -> Line Chart with Currency, Smoothing, Markers
        elif 'benchmarking' in name_lower:
             # Y-Axis (Price) -> Currency
             if chart.value_axis:
                 chart.value_axis.tick_labels.number_format = f'"{currency_symbol}" #,##0'
                 if currency_symbol == '$':
                     chart.value_axis.major_unit = 5000000
             
             # Lines: Smooth + Markers (for single point visibility)
             try:
                 for series in chart.series:
                     series.smooth = True
                     series.marker.style = XL_MARKER_STYLE.CIRCLE
                     series.marker.size = 7 # Visible but not huge
             except Exception as e:
                 print(f"Error setting markers: {e}")

        # Apply Brand Colors (if not varying by category)
        val_axis = chart.value_axis
        val_axis.has_major_gridlines = False # Clean Layout (No Grid)
        
        # Standard Color Cycle if not vary_by_categories
        if not plot.vary_by_categories and chart_type not in ['scatter']:
             brand_palette = [
                RGBColor(30, 41, 59),   # Dark Blue
                RGBColor(71, 85, 105),  # Light Blue
                RGBColor(59, 130, 246), # Accent Blue
                RGBColor(148, 163, 184) # Lighter Slate
            ]
             for i, series in enumerate(chart.series):
                color = brand_palette[i % len(brand_palette)]
                try:
                    if chart_type == 'line':
                         series.format.line.solid()
                         series.format.line.color.rgb = color
                    else:
                         series.format.fill.solid()
                         series.format.fill.fore_color.rgb = color
                except:
                    pass

    except Exception as e:
        print(f"Chart formatting warning: {e}")

    # Always add the table afterwards
    add_table_slide(prs, chart_info.get('chart_title', 'Datos'), rows, currency_symbol)
