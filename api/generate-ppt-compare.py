from http.server import BaseHTTPRequestHandler
import json
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.enum.text import PP_ALIGN
import sys
import os

# Import shared components
try:
    from api.ppt_shared import (
        DARK_BLUE, DEEP_NAVY, LIGHT_BLUE, WHITE,
        format_value, set_font, create_logo_slide, create_intro_slide,
        add_chart_slide, add_table_slide
    )
except ImportError:
    # Handle running as script vs module
    sys.path.append(os.path.dirname(os.path.abspath(__file__)))
    from ppt_shared import (
        DARK_BLUE, DEEP_NAVY, LIGHT_BLUE, WHITE,
        format_value, set_font, create_logo_slide, create_intro_slide,
        add_chart_slide, add_table_slide
    )

class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        content_len = int(self.headers.get('Content-Length', 0))
        post_body = self.rfile.read(content_len)
        data = json.loads(post_body)
        
        prs = Presentation()
        
        # 1. Cover
        create_logo_slide(prs)
        
        # 2. Intro
        today = datetime.now().strftime("%d/%m/%Y")
        title = data.get('reportTitle', 'Reporte Comparativo')
        create_intro_slide(prs, title, today)
        
        # 3. Currency Context
        currency_symbol = data.get('currency', '$')
        
        # 4. Content Slides (Charts & Tables)
        # Note: Compare page sends 'charts' and 'tables'
        slides_content = data.get('slides', [])
        
        # If slides are passed as a list
        for slide_data in slides_content:
            if slide_data.get('type') == 'chart':
                add_chart_slide(prs, slide_data, currency_symbol)
            elif slide_data.get('type') == 'table':
                add_table_slide(prs, slide_data.get('title', 'Tabla'), slide_data.get('data', []), currency_symbol)
                
        # 5. Closing Slide
        create_logo_slide(prs)
        
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        self.send_response(200)
        self.send_header('Content-type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        self.send_header('Content-Disposition', 'attachment; filename="reporte_comparar.pptx"')
        self.end_headers()
        self.wfile.write(ppt_stream.getvalue())
