from http.server import BaseHTTPRequestHandler
import json
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches
import sys
import os

# Import shared components
try:
    from api.ppt_shared import (
        create_logo_slide, create_intro_slide,
        add_chart_slide, add_table_slide
    )
except ImportError:
    sys.path.append(os.path.dirname(os.path.abspath(__file__)))
    from ppt_shared import (
        create_logo_slide, create_intro_slide,
        add_chart_slide, add_table_slide
    )

class handler(BaseHTTPRequestHandler):
    def do_POST(self):
        content_len = int(self.headers.get('Content-Length', 0))
        post_body = self.rfile.read(content_len)
        data = json.loads(post_body)
        
        prs = Presentation()
        # Enforce 16:9 Aspect Ratio (Widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        # 1. Cover
        create_logo_slide(prs)
        
        # 2. Intro
        today = datetime.now().strftime("%d/%m/%Y")
        title = data.get('reportTitle', 'Evolución de Precios')
        create_intro_slide(prs, title, today)
        
        # 3. Currency Context
        currency_symbol = data.get('currency', '$')
        
        # 4. Content Slides
        slides_content = data.get('slides', [])
        
        for slide_data in slides_content:
            if slide_data.get('type') == 'chart':
                add_chart_slide(prs, slide_data, currency_symbol)
            elif slide_data.get('type') == 'table':
                 add_table_slide(prs, slide_data.get('chart_title', 'Datos'), slide_data.get('data', []), currency_symbol)
                
        # 5. Closing
        create_logo_slide(prs)
        
        ppt_stream = io.BytesIO()
        prs.save(ppt_stream)
        ppt_stream.seek(0)
        
        self.send_response(200)
        self.send_header('Content-type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
        self.send_header('Content-Disposition', 'attachment; filename="reporte_evolucion.pptx"')
        self.end_headers()
        self.wfile.write(ppt_stream.getvalue())
