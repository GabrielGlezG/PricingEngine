from http.server import BaseHTTPRequestHandler
import json
import io
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData, XyChartData, BubbleChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import base64
import io
import sys

import sys
import os

# Ensure current directory is in path for sibling imports
current_dir = os.path.dirname(os.path.abspath(__file__))
if current_dir not in sys.path:
    sys.path.append(current_dir)

# Import shared components
try:
    from api.ppt_shared import (
        DARK_BLUE, DEEP_NAVY, LIGHT_BLUE, WHITE,
        format_value, set_font, get_image_stream,
        create_logo_slide, create_intro_slide,
        add_chart_slide, add_table_slide,
        LOGO_B64, BG_B64
    )
except ImportError:
    # Handle direct script execution where api package might not be resolved
    from ppt_shared import (
        DARK_BLUE, DEEP_NAVY, LIGHT_BLUE, WHITE,
        format_value, set_font, get_image_stream,
        create_logo_slide, create_intro_slide,
        add_chart_slide, add_table_slide,
        LOGO_B64, BG_B64
    )

def create_title_slide(prs, title, date_str):
    """Fallback title slide if no images available"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = date_str

def create_summary_slide(prs, summary, currency_symbol):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Resumen Ejecutivo"
    set_font(slide.shapes.title, font_name="Avenir Black", font_size=Pt(32), bold=True, color=DARK_BLUE)
    
    rows = 10
    cols = 2
    # 16:9 Layout: Center Summary Table
    # Width 8", Margin (13.33 - 8)/2 = 2.665"
    shape = slide.shapes.add_table(rows, cols, Inches(2.665), Inches(2), Inches(8), Inches(4))
    table = shape.table
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(4)
    
    # Header
    table.cell(0, 0).text = "Métrica"
    table.cell(0, 1).text = "Valor"
    for c in range(2):
        cell = table.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE # #1E293B
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.name = "Avenir Medium" # Template Header Font
    
    metrics = [
        ("Total Modelos", summary.get('total_models', 0), 'integer'),
        ("Total Marcas", summary.get('total_brands', 0), 'integer'),
        ("Precio Promedio", summary.get('avg_price', 0), 'currency'),
        ("Precio Mediano", summary.get('median_price', 0), 'currency'),
        ("Precio Mínimo", summary.get('min_price', 0), 'currency'),
        ("Precio Máximo", summary.get('max_price', 0), 'currency'),
        ("Desviación Estándar", summary.get('price_std_dev', 0), 'currency'),
        ("Coef. Variación", summary.get('variation_coefficient', 0), 'percent'),
        ("Descuento Promedio", summary.get('avg_discount_pct', 0), 'percent')
    ]
    
    for i, (label, val, fmt) in enumerate(metrics):
        row = i + 1
        c1 = table.cell(row, 0)
        c1.text = label
        c1.text_frame.paragraphs[0].font.name = "Avenir Medium"
        c2 = table.cell(row, 1)
        c2.text = format_value(val, fmt, currency_symbol)
        c2.text_frame.paragraphs[0].font.name = "Avenir Medium"
        c2.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT

def add_table_slide(prs, title, rows, currency_symbol='$'):
    if not rows: return
    MAX_ROWS = 12
    chunks = [rows[i:i + MAX_ROWS] for i in range(0, len(rows), MAX_ROWS)]
    
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide_title = f"{title}" if i == 0 else f"{title} (Cont.)"
        slide.shapes.title.text = slide_title
        set_font(slide.shapes.title, font_name="Avenir Black", font_size=Pt(28), bold=True, color=DARK_BLUE)
        
        headers = list(rows[0].keys())
        # 16:9 Layout: Width 12", Margin 0.665"
        shape = slide.shapes.add_table(len(chunk)+1, len(headers), Inches(0.665), Inches(1.5), Inches(12), Inches(0.4*(len(chunk)+1)))
        table = shape.table
        
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BLUE # #1E293B
            tf = cell.text_frame.paragraphs[0]
            tf.font.color.rgb = WHITE
            tf.font.bold = True
            tf.font.name = "Avenir Medium" # Template Header Font https://github.com/scanny/python-pptx/issues/455
            tf.font.size = Pt(10)
            tf.alignment = PP_ALIGN.CENTER

        for r, row_data in enumerate(chunk):
            for c, header in enumerate(headers):
                val = row_data.get(header)
                cell = table.cell(r + 1, c)
                
                h_str = str(header).lower()
                t_str = str(title).lower()
                fmt = None
                
                if any(x in h_str for x in ['%', 'percent', 'variación', 'variation', 'coef', 'descuento', 'volatilidad']):
                    fmt = 'percent'
                elif any(x in h_str for x in ['cantidad', 'cant.', 'volumen', 'versiones', 'total', 'numero', 'count']):
                    fmt = 'integer'
                elif any(x in h_str for x in ['precio', 'price', 'monto', 'valor', 'bono', 'lista', 'costo', 'avg', 'min', 'max', 'promedio']):
                    fmt = 'currency'
                elif any(x in t_str for x in ['volatilidad', 'volatility', 'tendencia', 'trend', 'variación', 'variation', 'share', 'participación', 'discount', 'descuento']):
                     if isinstance(val, (int, float)) and not any(x in h_str for x in ['fecha', 'date', 'year', 'año', 'mes']): 
                        fmt = 'percent'
                elif "precio" in t_str or "price" in t_str:
                     if isinstance(val, (int, float)): fmt = 'currency'
                
                if 'año' in h_str or 'year' in h_str: fmt = None

                cell.text = format_value(val, fmt, currency_symbol)
                tf = cell.text_frame.paragraphs[0]
                tf.font.name = "Avenir Medium" # Template Body Font
                tf.font.size = Pt(9)
                
                if fmt in ['currency', 'percent', 'integer']:
                     tf.alignment = PP_ALIGN.RIGHT
                else:
                     tf.alignment = PP_ALIGN.LEFT



def generate_ppt(data):
    try:
        prs = Presentation()
        # Enforce 16:9 Aspect Ratio (Widescreen)
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
        title = data.get('title', 'Reporte Dashboard')
        currency_symbol = data.get('currencySymbol', '$')
        date_str = datetime.now().strftime("%d/%m/%Y")
        
        # 1. Slide 1: Logo Cover
        create_logo_slide(prs)
    
        # 2. Slide 2: Intro
        create_intro_slide(prs, title, date_str)
        
        # 3. Summary Slide (Same as Excel Summary Sheet)
        summary = data.get('summary')
        if summary:
            try:
                create_summary_slide(prs, summary, currency_symbol)
            except Exception as e:
                print(f"Summary slide error: {e}")
            
        # 4. Sheets (Charts + Data Tables)
        sheets = data.get('sheets', [])
        for sheet in sheets:
            try:
                add_chart_slide(prs, sheet, currency_symbol)
            except Exception as e:
                print(f"Error creating chart/table slide {sheet.get('name')}: {e}")
                
        # 5. Models Data (Raw Table)
        models = data.get('models', [])
        if models:
            model_rows = []
            for m in models:
                 p_lista = float(m.get('precio_lista', 0) or 0)
                 bono = float(m.get('bono', 0) or 0)
                 dsc = (bono / p_lista) if p_lista > 0 else 0
                 
                 model_rows.append({
                     "Marca": m.get('brand'),
                     "Modelo": m.get('model'),
                     "Versión": m.get('submodel', '-'),
                     "Estado": m.get('estado', 'N/A'),
                     "Precio Lista": m.get('precio_lista', 0),
                     "Bono": m.get('bono', 0),
                     "Precio Final": m.get('precio_con_bono', 0),
                     "% Desc.": dsc
                 })
            
            try:
                add_table_slide(prs, "Detalle de Modelos", model_rows, currency_symbol)
            except Exception as e:
                print(f"Error adding models table: {e}")

        # 6. Last Slide: Logo Cover
        create_logo_slide(prs)
            
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        return output.getvalue()
    except Exception as global_e:
        # Fallback: Create a simple error presentation
        print(f"CRITICAL ERROR GENERATING PPT: {global_e}")
        import traceback
        traceback.print_exc()
        
        err_prs = Presentation()
        slide = err_prs.slides.add_slide(err_prs.slide_layouts[0])
        slide.shapes.title.text = "Error Generando Reporte"
        slide.placeholders[1].text = f"Detalles: {str(global_e)}\n\nConsulte los logs del servidor."
        
        output = io.BytesIO()
        err_prs.save(output)
        output.seek(0)
        return output.getvalue()

class handler(BaseHTTPRequestHandler):
    def do_OPTIONS(self):
        self.send_response(200)
        self.send_header('Access-Control-Allow-Origin', '*')
        self.send_header('Access-Control-Allow-Methods', 'POST, OPTIONS')
        self.send_header('Access-Control-Allow-Headers', 'Content-Type')
        self.end_headers()
    
    def do_POST(self):
        try:
            content_length = int(self.headers['Content-Length'])
            post_data = self.rfile.read(content_length)
            data = json.loads(post_data.decode('utf-8'))
            filename = data.get('filename', 'Presentation.pptx').replace('.xlsx', '.pptx')
            ppt_bytes = generate_ppt(data)
            self.send_response(200)
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Content-Type', 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
            self.send_header('Content-Disposition', f'attachment; filename="{filename}"')
            self.send_header('Content-Length', len(ppt_bytes))
            self.end_headers()
            self.wfile.write(ppt_bytes)
        except Exception as e:
            self.send_response(500)
            self.send_header('Access-Control-Allow-Origin', '*')
            self.send_header('Content-Type', 'application/json')
            self.end_headers()
            self.wfile.write(json.dumps({"error": str(e)}).encode())
