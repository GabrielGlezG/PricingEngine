from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys
import os

def hex_color(color_obj):
    if not color_obj: return "None"
    try:
        # Check type
        if hasattr(color_obj, 'type'):
             if str(color_obj.type) == 'RGB (1)':
                 return str(color_obj.rgb)
             elif str(color_obj.type) == 'THEME (2)':
                 return f"Theme({color_obj.theme_color})"
    except:
        pass
    return "Unknown"

def inspect_ppt(path):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return

    try:
        prs = Presentation(path)
        print(f"Inspecting: {path}")
        
        for i, slide in enumerate(prs.slides):
            print(f"\n--- Slide {i+1} ---")
            
            # Background
            if slide.background:
                 try:
                     fill = slide.background.fill
                     if fill.type == MSO_SHAPE_TYPE.PICTURE: 
                          print("  Background is Picture")
                     elif fill.fore_color:
                          print(f"  Background Color: {hex_color(fill.fore_color)}")
                 except:
                     pass

            if slide.shapes.title:
                t = slide.shapes.title
                print(f"  Title: '{t.text_frame.text.strip()}'")
                if t.text_frame.paragraphs:
                    p = t.text_frame.paragraphs[0]
                    run = p.runs[0] if p.runs else None
                    if run:
                        font = run.font
                        print(f"    Font: {font.name}, Size: {font.size}, Color: {hex_color(font.color)}")
            
            for shape in slide.shapes:
                if shape.has_table:
                    print("  [Found Table]")
                    tbl = shape.table
                    try:
                        c00 = tbl.cell(0, 0)
                        fill = c00.fill
                        print(f"    Header Fill: {hex_color(fill.fore_color)}")
                        if c00.text_frame.paragraphs:
                            if c00.text_frame.paragraphs[0].runs:
                                font = c00.text_frame.paragraphs[0].runs[0].font
                                print(f"    Header Font: {font.name}, Size: {font.size}, Color: {hex_color(font.color)}")
                            else:
                                # Check paragraph formatting if no runs
                                font = c00.text_frame.paragraphs[0].font
                                print(f"    Header Font (Para): {font.name}, Size: {font.size}, Color: {hex_color(font.color)}")

                    except Exception as e:
                        print(f"    Error reading table: {e}")
                
                if shape.has_chart:
                    title = "None"
                    try:
                        title = shape.chart.chart_title.text_frame.text
                    except: pass
                    print(f"  [Found Chart] Title: {title}")

    except Exception as e:
        print(f"CRITICAL ERROR: {e}")

inspect_ppt("public/FORMATO MAURICIO.pptx")
